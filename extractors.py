"""
内容提取模块 ─ 为不同类型文件实现内容抽取。

所有提取函数签名统一为:
    extract_xxx(path: Path, max_chars: int) -> str

提取失败时返回空字符串或错误提示（不抛异常）。
"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Callable, Dict, List

from config import AgentConfig
from scanner import FileInfo

# --------------- 提取器注册表 ---------------

_EXTRACTORS: Dict[str, Callable[[Path, int], str]] = {}


def _register(*extensions: str):
    """装饰器：按扩展名注册提取函数。"""
    def decorator(func: Callable[[Path, int], str]):
        for ext in extensions:
            _EXTRACTORS[ext.lower()] = func
        return func
    return decorator


# --------------- 各类提取器 ---------------

# ---- 代码 / 纯文本 ----
_CODE_EXTS = (
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".c", ".cpp", ".h", ".hpp",
    ".go", ".rs", ".rb", ".php", ".sh",
    ".bat", ".ps1", ".sql", ".r", ".m",
    ".css", ".html", ".xml", ".json", ".yaml", ".yml",
    ".md", ".txt", ".csv", ".log", ".ini", ".toml", ".cfg",
)


@_register(*_CODE_EXTS)
def _extract_text(path: Path, max_chars: int) -> str:
    """读取文本文件并提取摘要。对代码文件额外抓取函数定义行。"""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[读取失败: {e}]"

    # 对 Python / JS / Go 等代码文件提取函数签名
    ext = path.suffix.lower()
    if ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs", ".cpp", ".c", ".h", ".hpp"):
        patterns = [
            r"^\s*(def |class |async def )",        # Python
            r"^\s*(function |const \w+ ?= ?\()",    # JS/TS
            r"^\s*func ",                            # Go
            r"^\s*(fn |pub fn )",                    # Rust
            r"^\s*(public |private |protected )",    # Java/C++
        ]
        combined = "|".join(f"({p})" for p in patterns)
        sig_lines = [line.rstrip() for line in text.splitlines() if re.search(combined, line)]
        if sig_lines:
            sigs = "\n".join(sig_lines[:20])
            return f"[函数签名]\n{sigs}\n\n[前文摘要]\n{text[:max_chars]}"

    return text[:max_chars]


# ---- Word (.docx) ----
@_register(".docx")
def _extract_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)[:max_chars]
    except ImportError:
        return "[未安装 python-docx，无法解析 .docx]"
    except Exception as e:
        return f"[解析 .docx 失败: {e}]"


# ---- PPT (.pptx) ----
@_register(".pptx")
def _extract_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        text_parts: List[str] = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                text_parts.append(f"[幻灯片{slide_idx}] " + " | ".join(slide_texts))
        return "\n".join(text_parts)[:max_chars]
    except ImportError:
        return "[未安装 python-pptx，无法解析 .pptx]"
    except Exception as e:
        return f"[解析 .pptx 失败: {e}]"


# ---- PPT 旧格式 (.ppt) ----
@_register(".ppt")
def _extract_ppt_legacy(path: Path, max_chars: int) -> str:
    return "[.ppt 为旧版 PowerPoint 格式，暂不支持解析，建议转换为 .pptx]"


# ---- Excel (.xlsx / .xls) ----
@_register(".xlsx", ".xls")
def _extract_excel(path: Path, max_chars: int) -> str:
    try:
        import pandas as pd
        # 读取所有 sheet 的前 5 行
        xls = pd.ExcelFile(str(path))
        parts: List[str] = []
        for sheet_name in xls.sheet_names[:5]:  # 最多读 5 个 sheet
            df = pd.read_excel(xls, sheet_name=sheet_name, nrows=5)
            parts.append(f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}")
        return "\n\n".join(parts)[:max_chars]
    except ImportError:
        return "[未安装 pandas/openpyxl，无法解析 Excel]"
    except Exception as e:
        return f"[解析 Excel 失败: {e}]"


# ---- PDF ----
@_register(".pdf")
def _extract_pdf(path: Path, max_chars: int) -> str:
    try:
        import pdfplumber
        text_parts: List[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:3]:  # 前 3 页
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)[:max_chars]
    except ImportError:
        return "[未安装 pdfplumber，无法解析 PDF]"
    except Exception as e:
        return f"[解析 PDF 失败: {e}]"


# --------------- 对外接口 ---------------

def extract_content(fi: FileInfo, cfg: AgentConfig) -> str:
    """根据文件扩展名选取合适的提取器，返回内容摘要。"""
    ext = fi.ext.lower()
    extractor = _EXTRACTORS.get(ext)

    if extractor:
        return extractor(fi.path, cfg.max_content_chars)

    # 未注册的扩展名：尝试当作文本读取
    try:
        text = fi.path.read_text(encoding="utf-8", errors="strict")
        return text[:cfg.max_content_chars]
    except Exception:
        return ""


def enrich_file_list(files: List[FileInfo], cfg: AgentConfig) -> List[FileInfo]:
    """批量为 FileInfo 列表填充 content_summary 字段。"""
    for fi in files:
        fi.content_summary = extract_content(fi, cfg)
    return files
